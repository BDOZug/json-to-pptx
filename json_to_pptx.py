import json
from pptx import Presentation

def json_to_pptx(json_file, pptx_file):
    # JSON-Datei öffnen und einlesen
    with open(json_file, encoding='utf-8') as f:
        slides = json.load(f)

    # Neue Präsentation erzeugen
    prs = Presentation()
    # Layouts
    title_slide_layout = prs.slide_layouts[0]   # Titelfolie
    content_slide_layout = prs.slide_layouts[1] # Inhaltsfolie

    for i, slide in enumerate(slides):
        if i == 0:
            # Erste Folie als Titelfolie
            ppt_slide = prs.slides.add_slide(title_slide_layout)
            ppt_slide.shapes.title.text = slide.get('title', '')
            ppt_slide.placeholders[1].text = slide.get('content', '')
        else:
            # Danach Inhaltsfolien
            ppt_slide = prs.slides.add_slide(content_slide_layout)
            ppt_slide.shapes.title.text = slide.get('title', '')
            ppt_slide.placeholders[1].text = slide.get('content', '')

    # Präsentation speichern
    prs.save(pptx_file)
    print(f"Präsentation wurde erstellt: {pptx_file}")

if __name__ == "__main__":
    # Hier kannst du Dateinamen anpassen
    json_to_pptx("presentation.json", "output_presentation.pptx")
